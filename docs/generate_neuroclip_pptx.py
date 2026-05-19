"""
Generate NeuroClip FYP presentation (.pptx).
Run: pip install python-pptx pillow && python docs/generate_neuroclip_pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
DIAGRAMS = REPO_ROOT / "docs" / "diagrams_academic"
OUTPUT = REPO_ROOT / "docs" / "NeuroClip_FYP_Presentation.pptx"

EMERALD = RGBColor(0x10, 0xB9, 0x81)
SLATE = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def add_header_bar(slide, title: str, slide_num: int, total: int) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = SLATE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "NeuroClip"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.LEFT
    p2 = tf.add_paragraph()
    p2.text = f"{slide_num} / {total}"
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.RIGHT
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.55), SLIDE_W, Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = EMERALD
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SLATE


def add_bullets(slide, items: list[str], top: float = 1.55, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = BODY
        p.space_after = Pt(10)


def add_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    notes: str,
    slide_num: int,
    total: int,
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    add_header_bar(slide, title, slide_num, total)
    add_bullets(slide, bullets)
    add_notes(slide, notes)


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str,
    notes: str,
    slide_num: int,
    total: int,
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    add_header_bar(slide, title, slide_num, total)
    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path), Inches(0.6), Inches(1.45), width=Inches(12.1)
        )
    else:
        ph = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9.3), Inches(4)
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = LIGHT_BG
        ph.line.color.rgb = MUTED
        ph.text_frame.text = f"[Diagram: {image_path.name}]"
    cap = slide.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12), Inches(0.4))
    p = cap.text_frame.paragraphs[0]
    p.text = caption
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER
    add_notes(slide, notes)


def build_presentation() -> Presentation:
    prs = Presentation()
    set_slide_size(prs)
    total = 22

    # 1 Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = SLATE
    rect.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.9), SLIDE_W, Inches(0.6)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = EMERALD
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = "NeuroClip"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.0))
    tf = sub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Context-Aware Multimodal Video Processing\nfor Semantic Clip Retrieval"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tag = slide.shapes.add_textbox(Inches(1.2), Inches(4.1), Inches(10.9), Inches(0.6))
    p = tag.text_frame.paragraphs[0]
    p.text = "AI-powered semantic clip retrieval, privacy blurring, and smart compression"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.6))
    tf = meta.text_frame
    lines = [
        "Ali Javed  •  Asad Sardar",
        "Supervisor: Subhan Arif",
        "Department of Computer Science",
        "National University of Modern Languages — Faisalabad",
        "Final Year Project  •  2025–2026",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

    add_notes(
        slide,
        "Open with the problem: long videos are everywhere but finding one moment is painful. "
        "Introduce NeuroClip as your FYP solution. Transition: outline what you will cover.",
    )

    slides_data = [
        (
            "Agenda",
            [
                "Problem & motivation",
                "Objectives and scope",
                "System architecture & technology stack",
                "Processing pipeline (ASR + OCR + embeddings)",
                "Semantic retrieval & clip generation",
                "Privacy blurring & compression modules",
                "Evaluation, comparison & future work",
                "Demo flow & Q&A",
            ],
            "Walk through the agenda in 20 seconds. Tell examiners you will show architecture, "
            "methodology, and a live-oriented demo path at the end.",
        ),
        (
            "Problem & Motivation",
            [
                "Long-form lectures and meetings are hard to navigate manually.",
                "Keyword search misses intent — synonyms and paraphrases fail.",
                "Hits are fragmented snippets without surrounding context.",
                "Many pipelines re-transcribe and re-embed on every new query.",
                "Result: wasted time scrubbing timelines instead of learning.",
            ],
            "Emphasize three failure modes from your paper: semantic gap, fragmentation, redundancy. "
            "Q: Why not YouTube chapters? A: They are manual and not query-driven.",
        ),
        (
            "Proposed Solution",
            [
                "NeuroClip: a full-stack web platform that turns video into a queryable knowledge base.",
                "Natural-language queries return coherent, playable clips — not raw timestamps.",
                "Multimodal fusion: speech (ASR) + on-screen text (OCR) in one index.",
                "Persistent job-level caching: embeddings computed once, searched many times.",
            ],
            "One-liner: Upload once, ask in plain English, get the right clip. "
            "Transition to formal objectives next.",
        ),
        (
            "Project Objectives",
            [
                "Fuse ASR transcripts and frame OCR into searchable sentence-level units.",
                "Implement intent-aware retrieval with vector search, re-ranking, and FFmpeg clips.",
                "Provide query-driven privacy redaction (faces, IDs, license plates).",
                "Support codec-aware compression for lightweight sharing.",
                "Evaluate retrieval with grounded queries (Precision@K, nDCG@K, latency).",
                "Deploy a secure PWA with auth, history, and reusable job_id caching.",
            ],
            "Map each objective to a module you built. Examiners like traceability from objective to feature.",
        ),
        (
            "Scope & Target Users",
            [
                "In scope: upload/URL ingestion, transcription, embeddings, semantic search, clips.",
                "In scope: blurring, compression, Supabase auth, processing history.",
                "Users: students, educators, corporate teams reviewing lectures and meetings.",
                "Environment: web browser (PWA); backend on FastAPI; cloud-ready storage.",
                "Out of scope (future): full visual-only search without text channel.",
            ],
            "Clarify boundaries — shows maturity. Mention educational video focus for evaluation.",
        ),
    ]

    n = 2
    for title, bullets, notes in slides_data:
        add_content_slide(prs, title, bullets, notes, n, total)
        n += 1

    add_image_slide(
        prs,
        "System Architecture",
        DIAGRAMS / "4.6_system_sequence.png",
        "Figure: End-to-end system sequence (User → Frontend → Backend → Supabase → Clips)",
        "Explain actors: user, React frontend, FastAPI, AssemblyAI, EasyOCR, pgvector, FFmpeg. "
        "This is your Figure 1 equivalent from the research paper.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Technology Stack",
        [
            "Frontend: React 18, TypeScript, Vite, Tailwind CSS, shadcn/ui, TanStack Query",
            "Backend: FastAPI (Python 3.11+), PyTorch, sentence-transformers",
            "AI services: AssemblyAI (ASR), EasyOCR (visual text), YOLOv8 (detection)",
            "Media: FFmpeg (clip/blur/compress), yt-dlp (URL download)",
            "Data: Supabase — PostgreSQL, pgvector, Auth, Storage",
            "Deployment: Progressive Web App; Docker compose available",
        ],
        "Keep this slide fast. Highlight why Supabase: auth + vectors + storage in one stack.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Methodology Overview",
        [
            "1. Requirements analysis — use cases, functional/non-functional needs",
            "2. System design — DFDs, sequence diagrams, component & deployment models",
            "3. Implementation — modular FastAPI pipeline + React PWA",
            "4. Testing — 50 structured queries; compression profile sweeps",
            "5. Deployment — Supabase-backed production configuration",
        ],
        "Align with SDLC chapter. Point to diagrams in your report (Ch. 4).",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Data Processing Pipeline",
        [
            "Upload video or URL → assign UUID job_id",
            "AssemblyAI: audio → SRT → sentence units (text, start, end)",
            "EasyOCR: sample frames every ~3s → merge text into nearest sentences",
            "Embed each enriched sentence with all-MiniLM-L6-v2 (384 dimensions)",
            "Persist to Supabase pgvector + local JSON artifacts",
            "User query → retrieve windows → merge → FFmpeg playable clips",
        ],
        "Walk left-to-right: Upload → ASR+OCR → Embed → Store → Search → Clips. "
        "Matches your standee pipeline badges.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Multimodal Fusion (ASR + OCR)",
        [
            "ASR captures spoken explanations and dialogue.",
            "OCR captures slides, code on screen, whiteboard equations, captions.",
            "OCR snippets merged by timestamp into nearest transcript sentence.",
            "Unmatched OCR becomes synthetic visual sentences for indexing.",
            "Enables queries like “slide about hash collision” without exact spoken words.",
        ],
        "This is a key differentiator vs transcript-only tools. Give a concrete lecture example.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Semantic Retrieval",
        [
            "Query encoded in same 384-dim space as sentence embeddings.",
            "Cosine similarity via pgvector (fallback when LLM route unavailable).",
            "Windowed retrieval: groups of W consecutive sentences (default W=5, stride=2).",
            "Optional cross-encoder re-ranking (ms-marco-MiniLM-L-6-v2) on top candidates.",
            "Neighbor merging joins adjacent high-score windows into one coherent clip.",
        ],
        "Explain why windows beat single sentences: context for the viewer. "
        "Mention graceful fallback if re-ranker fails.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Clip Extraction & Playback",
        [
            "Merged time ranges sent to FFmpeg with configurable padding (±1.5s default).",
            "Stream-copy when possible; re-encode when blur/compress applied.",
            "Clips stored in Supabase Storage; URLs returned to React player.",
            "History keyed by job_id — revisit past searches without reprocessing.",
            "Configurable min/max clip duration (e.g. 10s–120s).",
        ],
        "Stress reproducibility: same job_id, same embeddings, faster repeat queries.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Privacy Blurring Module",
        [
            "Natural-language instruction: e.g. “blur faces between 2:30 and 4:00”.",
            "YOLOv8 detects faces, ID cards, license plates in sampled frames.",
            "Gaussian blur applied per bounding box; tracking reduces flicker.",
            "FFmpeg re-encodes redacted segment for download/sharing.",
            "Supports compliance-friendly clip sharing from sensitive recordings.",
        ],
        "Demo tip: show before/after on a short segment. Q: Real-time? A: Post-processing on demand.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Video Compression Module",
        [
            "POST /compress-video — distribution-ready outputs.",
            "Codec: H.265/HEVC — GPU hevc_nvenc or CPU libx265 fallback.",
            "Default cap 720p; CRF controls quality vs size trade-off.",
            "Response includes original size, compressed size, reduction %, time.",
            "Typical 60–78% size reduction on educational profiles (per evaluation).",
        ],
        "Pair with sharing use case: student sends clip over slow network.",
        n,
        total,
    )
    n += 1

    add_image_slide(
        prs,
        "Use Case Model",
        DIAGRAMS / "4.2_use_case.png",
        "Figure: Primary use cases — upload, search, history, blur, compress",
        "Primary actor: End User. Walk through upload → search → play → optional blur/compress.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Database & Caching",
        [
            "Supabase Auth: email/password, row-level security.",
            "Tables: videos, embeddings, sentence embeddings, processing history.",
            "pgvector: server-side nearest-neighbor search (match_video_embeddings RPC).",
            "job_id links file, transcript, vectors, and search history.",
            "Cached queries: sub-second latency when embeddings already exist.",
        ],
        "Explain why caching matters for FYP demo: second query on same video is instant.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Evaluation & Results",
        [
            "Benchmark: 50 structured queries — math, CS, AI, cybersecurity domains.",
            "Metrics: Precision@K, Recall@K, nDCG@K, temporal IoU, boundary error.",
            "Also measure end-to-end query latency and compression ratio / SSIM where available.",
            "Qualitative: semantically relevant clips with coherent context windows.",
            "Full pipeline completes in minutes on consumer hardware; cached search sub-second.",
        ],
        "Be honest: cite your eval pack in docs/evaluation. Do not invent precise P@5 numbers if not in thesis.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Comparison with Existing Systems",
        [
            "CapCut / editors: manual; no semantic search or auto privacy pipeline.",
            "Summarify / Notta: transcription-focused; limited custom upload + retrieval.",
            "Keyword transcript search: misses paraphrase and on-screen slide content.",
            "NeuroClip: unified ASR+OCR fusion, vector search, clips, blur, compress, cache.",
        ],
        "Use the comparison table from your paper verbally if asked for related work.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Limitations & Future Work",
        [
            "ASR/OCR quality degrades with noise, handwriting, or fast scene changes.",
            "Cross-encoder improves precision but adds latency vs embedding-only path.",
            "GPU optional — compression may be slower on CPU-only machines.",
            "Future: CLIP/OpenCLIP frame embeddings for visual search.",
            "Future: richer speaker diarization and multi-language support.",
        ],
        "Shows critical thinking. Examiners respect known limits.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Key Contributions",
        [
            "Multimodal fusion of ASR + OCR into enriched sentence units.",
            "Dense vector semantic search with windowed, mergeable clip output.",
            "Two-stage retrieval: fast bi-encoder + optional cross-encoder re-ranking.",
            "Job-keyed persistent caching eliminating redundant reprocessing.",
            "Integrated privacy blurring and H.265 compression in one platform.",
            "Deployable PWA with Supabase auth and searchable history.",
        ],
        "Summarize contributions as numbered list — mirrors paper Section 1.",
        n,
        total,
    )
    n += 1

    add_content_slide(
        prs,
        "Live Demo Flow",
        [
            "1. Login (Supabase Auth)",
            "2. Upload lecture video or paste YouTube URL",
            "3. Wait for processing (transcript + embeddings)",
            "4. Query: e.g. “explain backpropagation” or “hash collision slide”",
            "5. Play returned clip; open history by job_id",
            "6. Optional: blur faces segment → compress for sharing",
        ],
        "If live demo fails, have a pre-processed job_id ready. This slide is your backup script.",
        n,
        total,
    )
    n += 1

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = SLATE
    rect.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.8), Inches(5.3), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = EMERALD
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Questions?"
    p2.font.size = Pt(28)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    p3 = tf.add_paragraph()
    p3.text = "Ali Javed  •  Asad Sardar  |  Supervisor: Subhan Arif\nNUML Faisalabad — Computer Science"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(24)

    add_notes(slide, "Invite questions. Suggest demo if time permits. Have architecture diagram ready in report.")

    return prs


def main() -> None:
    prs = build_presentation()
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
